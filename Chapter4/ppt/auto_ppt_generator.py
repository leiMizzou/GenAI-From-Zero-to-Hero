from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import re

class PPTGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.set_theme()
        
    def set_theme(self):
        """设置专业配色方案"""
        self.theme_colors = {
            "primary": RGBColor(0, 112, 192),
            "secondary": RGBColor(155, 194, 230),
            "accent": RGBColor(79, 129, 189),
            "text": RGBColor(51, 51, 51)
        }
        
    def add_title_slide(self, title, subtitle=""):
        """添加标题幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = subtitle
        
        # 设置标题样式
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.theme_colors["primary"]
        
        # 设置副标题样式
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = self.theme_colors["accent"]
        
        return slide
    
    def add_content_slide(self, title, content, level=1):
        """添加内容幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = title
        content_shape.text = content
        
        # 设置标题样式
        title_shape.text_frame.paragraphs[0].font.size = Pt(32 if level == 1 else 28)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.theme_colors["primary"]
        
        # 设置内容样式
        for paragraph in content_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = self.theme_colors["text"]
            
        return slide
    
    def add_code_slide(self, title, code):
        """添加代码展示幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # 仅标题版式
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # 添加代码框
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        
        # 添加代码内容
        p = tf.add_paragraph()
        p.text = code
        p.font.name = 'Courier New'
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0, 0, 0)
        
        return slide
    
    def save(self, filename):
        """保存PPT文件"""
        self.prs.save(filename)
        print(f"PPT已保存为 {filename}")

def analyze_code(filepath):
    """分析代码文件提取关键信息"""
    with open(filepath, 'r', encoding='utf-8') as f:
        code = f.read()
    
    # 提取类定义
    class_pattern = r'class\s+(\w+).*?def\s+__init__\(.*?\):(.*?)(?=def\s+\w+|\Z)'
    classes = re.findall(class_pattern, code, re.DOTALL)
    
    # 提取函数定义
    func_pattern = r'def\s+(\w+)\(.*?\):(.*?)(?=def\s+\w+|\Z)'
    functions = re.findall(func_pattern, code, re.DOTALL)
    
    return {
        'classes': classes,
        'functions': functions,
        'imports': [line for line in code.split('\n') if line.startswith('import') or line.startswith('from')]
    }

def generate_ppt(code_file, output_file):
    """生成PPT主函数"""
    # 分析代码
    analysis = analyze_code(code_file)
    
    # 创建PPT生成器
    ppt = PPTGenerator()
    
    # 1. 添加封面
    ppt.add_title_slide(
        "长文生成器技术解析", 
        "基于两阶段架构的AI内容生成系统"
    )
    
    # 2. 添加架构概述
    ppt.add_content_slide(
        "系统架构",
        "• 两阶段处理流程\n" +
        "  1. 大纲生成阶段\n" +
        "  2. 内容扩展阶段\n\n" +
        "• 核心组件:\n" +
        "  - OutlineGenerator: 大纲生成器\n" +
        "  - ContentExpander: 内容扩展器\n" +
        "  - ArticleFormatter: 文章格式化"
    )
    
    # 3. 添加类定义
    for class_name, class_body in analysis['classes']:
        ppt.add_content_slide(
            f"类: {class_name}",
            f"功能描述:\n{class_body.strip()}\n\n" +
            "主要方法:\n" +
            "\n".join([f"- {func[0]}" for func in analysis['functions'] if func[0].startswith(class_name.lower())])
        )
    
    # 4. 添加核心方法
    ppt.add_content_slide(
        "核心方法",
        "\n".join([f"• {func[0]}: {func[1].strip()[:100]}..." for func in analysis['functions']])
    )
    
    # 5. 添加代码示例
    ppt.add_code_slide(
        "主函数代码",
        "\n".join([line for line in open(code_file, 'r').readlines() if not line.strip().startswith('#')][:30])
    )
    
    # 6. 添加使用示例
    ppt.add_content_slide(
        "使用示例",
        "命令行参数:\n" +
        "• --topic: 文章主题\n" +
        "• --style: 文章风格\n" +
        "• --length: 文章长度\n" +
        "• --output: 输出文件\n" +
        "• --api_key: OpenRouter API密钥"
    )
    
    # 保存PPT
    ppt.save(output_file)

if __name__ == "__main__":
    # Use relative path from the project root
    input_file = "Chapter4/examples/longwriter/main.py" 
    # Save the output in the longwriter directory
    output_file = "Chapter4/examples/longwriter/longwriter.pptx" 
    generate_ppt(input_file, output_file)